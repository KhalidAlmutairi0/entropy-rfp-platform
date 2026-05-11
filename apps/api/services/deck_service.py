"""Deck generation pipeline: Upload (.pptx/.docx) → Docling → Structured JSON → Claude → PptxGenJS → .pptx

Pipeline stages (each raises RuntimeError with a clear prefix on failure):
  Stage 1 — Parse:   Docling (preferred) or python-pptx/python-docx fallback
  Stage 2 — Claude:  claude-opus-4-6 generates slides from structured JSON + RFP content
  Stage 3 — Render:  Node.js PptxGenJS renders the .pptx from Claude's slide JSON
"""

import asyncio
import io
import json
import os
import tempfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

import structlog

from core.config import settings

logger = structlog.get_logger()

# Thread pool for running synchronous Docling / python-pptx calls without blocking the event loop
_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="docling")

_RENDERER_PATH = Path(__file__).parent.parent.parent / "deck-renderer" / "index.js"

# ── Structured JSON schema (shared between stages) ────────────────────────────
#
# {
#   "source_type": "pptx" | "docx",
#   "slides": [
#     {
#       "index": 0,
#       "layout": "title" | "content" | "two_column" | "image" | "blank",
#       "title": "...",
#       "body": ["bullet 1", "bullet 2"],
#       "notes": "...",
#       "theme": { "bg": "#ffffff", "accent": "#0070f3" }
#     }
#   ],
#   "global_theme": { "bg": "#ffffff", "accent": "#0070f3", "font": "Arial" }
# }


# ─────────────────────────────────────────────────────────────────────────────
# Stage 1 — Parse template with Docling (or fallback)
# ─────────────────────────────────────────────────────────────────────────────

async def parse_template(file_bytes: bytes, filename: str) -> dict:
    """Parse a .pptx or .docx template file into structured JSON.

    For .pptx files: uses python-pptx (Docling does not support PPTX).
    For .docx files: tries Docling first, falls back to python-docx with a warning logged.

    Raises:
        ValueError: unsupported file extension.
        RuntimeError: parsing failed at all levels.
    """
    suffix = Path(filename).suffix.lower()
    loop = asyncio.get_event_loop()

    if suffix == ".pptx":
        try:
            return await loop.run_in_executor(_executor, _parse_pptx_sync, file_bytes)
        except Exception as exc:
            raise RuntimeError(f"PPTX parsing failed: {exc}") from exc

    if suffix in (".docx", ".doc"):
        # 1. Try Docling
        try:
            result = await loop.run_in_executor(_executor, _parse_docx_with_docling_sync, file_bytes)
            return result
        except ImportError:
            logger.warning("deck_service: Docling not installed — falling back to python-docx")
        except Exception as exc:
            logger.warning("deck_service: Docling parse failed — falling back to python-docx", error=str(exc))

        # 2. Fallback: python-docx
        try:
            return await loop.run_in_executor(_executor, _parse_docx_fallback_sync, file_bytes)
        except Exception as exc:
            raise RuntimeError(f"DOCX parsing failed (Docling and python-docx both failed): {exc}") from exc

    raise ValueError(f"Unsupported template file type: '{suffix}'. Supported: .pptx, .docx")


# ── python-pptx parser ────────────────────────────────────────────────────────

def _parse_pptx_sync(file_bytes: bytes) -> dict:
    from pptx import Presentation  # type: ignore[import]

    prs = Presentation(io.BytesIO(file_bytes))
    global_theme = _extract_pptx_theme(prs)
    slides: list[dict] = []

    for idx, slide in enumerate(prs.slides):
        slide_data: dict = {
            "index": idx,
            "layout": _infer_pptx_layout(slide, idx),
            "title": "",
            "body": [],
            "notes": "",
            "theme": dict(global_theme),
        }

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = getattr(shape, "placeholder_format", None)
            paragraphs = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]

            if ph is not None:
                if ph.idx == 0:  # Title placeholder (idx 0 is always the title)
                    slide_data["title"] = shape.text_frame.text.strip()
                else:
                    slide_data["body"].extend(paragraphs)
            elif "title" in shape.name.lower():
                slide_data["title"] = slide_data["title"] or shape.text_frame.text.strip()
            else:
                slide_data["body"].extend(paragraphs)

        if slide.has_notes_slide:
            slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        slides.append(slide_data)

    return {"source_type": "pptx", "slides": slides, "global_theme": global_theme}


def _extract_pptx_theme(prs) -> dict:
    theme: dict = {"bg": "#FFFFFF", "accent": "#0070F3", "font": "Arial"}
    try:
        # Slide background from first slide
        if prs.slides:
            fill = prs.slides[0].background.fill
            if fill.type == 1:  # PP_FILL.SOLID
                theme["bg"] = f"#{fill.fore_color.rgb}"
    except Exception:
        pass
    return theme


def _infer_pptx_layout(slide, idx: int) -> str:
    layout_name = (slide.slide_layout.name or "").lower() if slide.slide_layout else ""
    if idx == 0 or "title" in layout_name and "content" not in layout_name:
        return "title"
    if "two" in layout_name or "comparison" in layout_name or "column" in layout_name:
        return "two_column"
    if "blank" in layout_name:
        return "blank"
    if "image" in layout_name or "picture" in layout_name:
        return "image"
    return "content"


# ── Docling DOCX parser ───────────────────────────────────────────────────────

def _parse_docx_with_docling_sync(file_bytes: bytes) -> dict:
    from docling.document_converter import DocumentConverter  # type: ignore[import]

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        f.write(file_bytes)
        tmp_path = f.name

    try:
        converter = DocumentConverter()
        result = converter.convert(tmp_path)
        doc = result.document

        slides: list[dict] = []
        current: dict | None = None
        slide_idx = 0

        for item, _ in doc.iterate_items():
            label = str(getattr(item, "label", "")).lower()
            text = getattr(item, "text", "").strip()
            if not text:
                continue

            is_heading = "heading" in label or "title" in label or "section" in label
            if is_heading:
                if current:
                    slides.append(current)
                current = {
                    "index": slide_idx,
                    "layout": "title" if slide_idx == 0 else "content",
                    "title": text,
                    "body": [],
                    "notes": "",
                    "theme": {},
                }
                slide_idx += 1
            elif current:
                current["body"].append(text)
            else:
                # Content before the first heading — start an untitled slide
                current = {
                    "index": slide_idx,
                    "layout": "content",
                    "title": "",
                    "body": [text],
                    "notes": "",
                    "theme": {},
                }
                slide_idx += 1

        if current:
            slides.append(current)

        return {"source_type": "docx", "slides": slides, "global_theme": {}}
    finally:
        os.unlink(tmp_path)


# ── python-docx fallback parser ───────────────────────────────────────────────

def _parse_docx_fallback_sync(file_bytes: bytes) -> dict:
    from docx import Document  # type: ignore[import]

    doc = Document(io.BytesIO(file_bytes))
    slides: list[dict] = []
    current: dict | None = None
    slide_idx = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        is_heading = "heading" in para.style.name.lower()
        if is_heading:
            if current:
                slides.append(current)
            current = {
                "index": slide_idx,
                "layout": "title" if slide_idx == 0 else "content",
                "title": text,
                "body": [],
                "notes": "",
                "theme": {},
            }
            slide_idx += 1
        elif current:
            current["body"].append(text)
        else:
            current = {
                "index": slide_idx,
                "layout": "content",
                "title": "",
                "body": [text],
                "notes": "",
                "theme": {},
            }
            slide_idx += 1

    if current:
        slides.append(current)

    return {"source_type": "docx", "slides": slides, "global_theme": {}}


# ─────────────────────────────────────────────────────────────────────────────
# Stage 2 — Claude generates slides
# ─────────────────────────────────────────────────────────────────────────────

_CLAUDE_SYSTEM_PROMPT = """You are an expert Arabic/English presentation designer for Entropy, a Saudi AI technology company.

You receive:
1. A structured JSON that describes a TEMPLATE presentation — its slide layouts, visual patterns, and content structure.
2. The RFP (tender) text that defines the content for the new presentation.

Your task: generate a NEW presentation for this RFP, following the SAME slide layout patterns from the template.

STRICT rules:
- Output ONLY a raw JSON array of slide objects. No markdown fences, no prose, no explanations.
- Use this exact per-slide schema:
  { "index": <int>, "layout": "<title|content|two_column|image|blank>", "title": "<string>", "body": ["<bullet>", ...], "notes": "<string>", "theme": { "bg": "<hex>", "accent": "<hex>" } }
- Generate between 8 and 15 slides. Required sections: Executive Summary, Understanding of Requirements, Proposed Methodology, Team & Expertise, Timeline, Pricing (mark as TBD), Why Entropy.
- Write content in the same primary language as the RFP. Default to Arabic.
- "body" is always an array of concise bullet-point strings (10–20 words each).
- Apply global_theme colors to every slide's "theme" field.
- Do NOT hallucinate. All claims must be traceable to the RFP text provided."""


async def call_claude_for_slides(structured_json: dict, rfp_text: str, rfp_title: str) -> list[dict]:
    """Send the structured template JSON + RFP content to Claude claude-opus-4-6.

    Returns a list of slide dicts matching the structured JSON schema.

    Raises:
        RuntimeError: Claude API error or invalid JSON response.
    """
    from services.llm_client import make_anthropic_client

    client = make_anthropic_client()

    template_summary = json.dumps(structured_json, ensure_ascii=False, indent=2)[:8000]
    rfp_excerpt = (rfp_text or "لا يوجد نص للمناقصة.")[:12000]
    global_theme_json = json.dumps(structured_json.get("global_theme", {}))

    user_message = (
        f"Template Presentation Structure:\n{template_summary}\n\n"
        f"Global Theme: {global_theme_json}\n\n"
        f"RFP Title: {rfp_title}\n\n"
        f"RFP Content:\n{rfp_excerpt}\n\n"
        "Generate the full presentation. Output raw JSON array only."
    )

    response = await client.messages.create(
        model="claude-opus-4-6",
        max_tokens=8192,
        system=_CLAUDE_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": user_message}],
    )

    raw = response.content[0].text.strip()

    # Strip markdown code fences if Claude wrapped the output
    if raw.startswith("```"):
        lines = raw.splitlines()
        raw = "\n".join(lines[1:])  # drop opening fence line
        if raw.endswith("```"):
            raw = raw[: raw.rfind("```")].strip()
    if raw.lstrip().startswith("json"):
        raw = raw.lstrip()[4:].strip()

    try:
        slides = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise RuntimeError(
            f"Claude returned invalid JSON. Decode error: {exc}. "
            f"First 500 chars of output: {raw[:500]!r}"
        ) from exc

    if not isinstance(slides, list):
        raise RuntimeError(
            f"Expected a JSON array from Claude, got {type(slides).__name__}. "
            f"Raw output (first 200 chars): {raw[:200]!r}"
        )

    return slides


# ─────────────────────────────────────────────────────────────────────────────
# Stage 3 — Render .pptx with PptxGenJS (Node.js subprocess)
# ─────────────────────────────────────────────────────────────────────────────

async def render_with_pptxgenjs(slides: list[dict], global_theme: dict) -> bytes:
    """Invoke the Node.js PptxGenJS renderer by writing JSON to its stdin.

    The renderer (apps/deck-renderer/index.js) reads JSON from stdin and writes
    the .pptx binary to stdout.

    Raises:
        FileNotFoundError: renderer script not found.
        RuntimeError: renderer process failed or timed out.
    """
    if not _RENDERER_PATH.exists():
        raise FileNotFoundError(
            f"PptxGenJS renderer not found at {_RENDERER_PATH}. "
            "Run: cd apps/deck-renderer && npm install"
        )

    payload = json.dumps(
        {"slides": slides, "global_theme": global_theme},
        ensure_ascii=False,
    ).encode("utf-8")

    proc = await asyncio.create_subprocess_exec(
        "node",
        str(_RENDERER_PATH),
        stdin=asyncio.subprocess.PIPE,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )

    try:
        stdout, stderr = await asyncio.wait_for(
            proc.communicate(input=payload),
            timeout=120.0,
        )
    except asyncio.TimeoutError as exc:
        proc.kill()
        await proc.wait()
        raise RuntimeError("PptxGenJS renderer timed out after 120 seconds") from exc

    if proc.returncode != 0:
        err_text = stderr.decode("utf-8", errors="replace")
        raise RuntimeError(
            f"PptxGenJS renderer exited with code {proc.returncode}: {err_text[:1000]}"
        )

    if not stdout:
        raise RuntimeError("PptxGenJS renderer produced empty output")

    return stdout


# ─────────────────────────────────────────────────────────────────────────────
# Full pipeline orchestrator
# ─────────────────────────────────────────────────────────────────────────────

async def generate_deck_pipeline(
    rfp_id: str,
    rfp_title: str,
    rfp_text: str,
    template_bytes: bytes,
    template_filename: str,
) -> bytes:
    """Run the full deck generation pipeline and return raw .pptx bytes.

    Stages:
      1. parse_template()       — Docling/python-pptx → structured JSON
      2. call_claude_for_slides() — Claude claude-opus-4-6 → slide JSON
      3. render_with_pptxgenjs() — Node.js PptxGenJS → .pptx bytes

    Each stage raises RuntimeError with a '[Stage N — Name]' prefix so the caller
    can surface per-stage error messages without inspecting internals.
    """
    # Stage 1
    logger.info("deck_pipeline: parsing template", rfp_id=rfp_id, filename=template_filename)
    try:
        structured_json = await parse_template(template_bytes, template_filename)
    except Exception as exc:
        raise RuntimeError(f"[Stage 1 — Parse] {exc}") from exc

    slide_count = len(structured_json.get("slides", []))
    logger.info("deck_pipeline: template parsed", rfp_id=rfp_id, slides=slide_count)

    # Stage 2
    logger.info("deck_pipeline: calling Claude", rfp_id=rfp_id, model="claude-opus-4-6")
    try:
        slides = await call_claude_for_slides(structured_json, rfp_text, rfp_title)
    except Exception as exc:
        raise RuntimeError(f"[Stage 2 — Claude] {exc}") from exc

    logger.info("deck_pipeline: slides generated", rfp_id=rfp_id, slides=len(slides))

    # Stage 3
    logger.info("deck_pipeline: rendering with PptxGenJS", rfp_id=rfp_id)
    try:
        pptx_bytes = await render_with_pptxgenjs(slides, structured_json.get("global_theme", {}))
    except Exception as exc:
        raise RuntimeError(f"[Stage 3 — Render] {exc}") from exc

    logger.info("deck_pipeline: done", rfp_id=rfp_id, size_bytes=len(pptx_bytes))
    return pptx_bytes


# ─────────────────────────────────────────────────────────────────────────────
# Background task wrapper (called by the FastAPI route via BackgroundTasks)
# ─────────────────────────────────────────────────────────────────────────────

async def _run_deck_pipeline_async(
    rfp_id: str,
    template_bytes: bytes,
    template_filename: str,
) -> None:
    """Background task: run the pipeline, store the .pptx in MinIO, update DB status."""
    from core.database import AsyncSessionLocal
    from models.rfp import RFP
    from services.storage import StorageService
    from sqlalchemy import select
    from sqlalchemy.orm import selectinload

    storage = StorageService.get_instance()

    async with AsyncSessionLocal() as db:
        result = await db.execute(
            select(RFP).options(selectinload(RFP.files)).where(RFP.id == rfp_id)
        )
        rfp = result.scalar_one_or_none()
        if not rfp:
            logger.error("deck_pipeline: RFP not found", rfp_id=rfp_id)
            return

        rfp.deck_status = "GENERATING"
        await db.commit()

        try:
            rfp_text = await _extract_rfp_text(rfp, storage)
            rfp_title = rfp.title_ar or rfp.title_en or "مناقصة"

            pptx_bytes = await generate_deck_pipeline(
                rfp_id=rfp_id,
                rfp_title=rfp_title,
                rfp_text=rfp_text,
                template_bytes=template_bytes,
                template_filename=template_filename,
            )

            pptx_path = f"decks/{rfp_id}/proposal.pptx"
            await storage.upload_file(
                pptx_bytes,
                pptx_path,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

            rfp.deck_pdf_path = pptx_path
            rfp.deck_status = "READY"
            await db.commit()
            logger.info("deck_pipeline: stored", rfp_id=rfp_id, path=pptx_path)

        except Exception as exc:
            logger.error("deck_pipeline: failed", rfp_id=rfp_id, error=str(exc))
            rfp.deck_status = "FAILED"
            await db.commit()


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

async def _load_template_from_storage(storage) -> tuple[bytes, str] | None:
    """Try to load a .pptx or .docx template from MinIO or the DECK_TEMPLATE_PATH env var.

    Returns (file_bytes, filename) or None if nothing is found.
    """
    # 1. MinIO — prefer .pptx, then .docx
    for minio_path, filename in [
        ("templates/master_proposal.pptx", "master_proposal.pptx"),
        ("templates/master_proposal.docx", "master_proposal.docx"),
    ]:
        try:
            content = await storage.download_file(minio_path)
            return content, filename
        except Exception:
            pass

    # 2. Local filesystem override via env var
    env_path = os.getenv("DECK_TEMPLATE_PATH", "")
    if env_path:
        p = Path(env_path)
        if p.exists() and p.suffix.lower() in (".pptx", ".docx"):
            return p.read_bytes(), p.name

    return None


async def _extract_rfp_text(rfp, storage) -> str:
    """Extract plain text from the RFP's uploaded files (reuses the ingestion task helper)."""
    from tasks.ingestion_tasks import _extract_text

    blocks: list[str] = []
    for rfp_file in rfp.files:
        try:
            content = await storage.download_file(rfp_file.storage_path)
            text = await _extract_text(content, rfp_file.mime_type or "", rfp_file.filename)
            if text.strip():
                blocks.append(f"=== {rfp_file.filename} ===\n{text.strip()[:10000]}")
        except Exception as exc:
            logger.warning("deck_pipeline: text extraction failed", file=rfp_file.filename, error=str(exc))

    return "\n\n".join(blocks)[:30000]
