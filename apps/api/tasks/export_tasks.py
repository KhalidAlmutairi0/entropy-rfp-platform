"""Celery tasks for proposal export (DOCX / PDF generation)."""

import asyncio
import io

import structlog

from core.celery_app import celery_app

logger = structlog.get_logger()


@celery_app.task(bind=True, name="tasks.export_tasks.export_proposal_task", max_retries=2)
def export_proposal_task(self, proposal_id: str, config: dict) -> dict:
    return asyncio.run(_export_async(proposal_id, config))


async def _export_async(proposal_id: str, config: dict) -> dict:
    from core.database import AsyncSessionLocal
    from models.proposal import Proposal
    from services.storage import StorageService
    from sqlalchemy import select
    from sqlalchemy.orm import selectinload

    storage = StorageService.get_instance()

    async with AsyncSessionLocal() as db:
        proposal = (
            await db.execute(
                select(Proposal).options(selectinload(Proposal.sections)).where(Proposal.id == proposal_id)
            )
        ).scalar_one_or_none()

        if not proposal:
            return {"error": "Proposal not found"}

        fmt = config.get("format", "pdf")
        language = config.get("language", "ar")
        results = {}

        if fmt in ("docx", "both"):
            docx_bytes = await _generate_docx(proposal, config, language)
            path = f"exports/{proposal_id}/proposal.docx"
            await storage.upload_file(docx_bytes, path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
            results["docx_path"] = path

        if fmt in ("pdf", "both"):
            docx_bytes = await _generate_docx(proposal, config, language)
            pdf_bytes = _docx_to_pdf(docx_bytes)
            path = f"exports/{proposal_id}/proposal.pdf"
            await storage.upload_file(pdf_bytes, path, "application/pdf")
            results["pdf_path"] = path

        logger.info("Proposal exported", proposal_id=proposal_id, format=fmt)
        return results


async def _generate_docx(proposal, config: dict, language: str) -> bytes:
    """Generate a DOCX file from proposal sections using python-docx."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Cover page
    if config.get("include_cover", True):
        cover = doc.add_paragraph()
        cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = cover.add_run("Entropy.sa\n")
        run.bold = True
        run.font.size = Pt(24)
        run.font.color.rgb = RGBColor(0x48, 0x65, 0x81)  # primary-600

        title = doc.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title.add_run("Technical Proposal\n")
        title_run.bold = True
        title_run.font.size = Pt(18)
        doc.add_page_break()

    # Table of contents placeholder
    if config.get("include_toc", True):
        toc_heading = doc.add_heading("Table of Contents", level=1)
        doc.add_paragraph("[Auto-generated TOC — update fields in Word]")
        doc.add_page_break()

    # Sections
    for idx, section in enumerate(proposal.sections, 1):
        title = section.title_en if language == "en" else (section.title_ar or section.title_en or f"Section {idx}")
        content = section.content_en if language == "en" else (section.content_ar or section.content_en or "")

        heading = doc.add_heading(f"{idx}. {title}" if config.get("include_section_numbers", True) else title, level=1)

        if content:
            for paragraph in content.split("\n\n"):
                if paragraph.strip():
                    doc.add_paragraph(paragraph.strip())

    # Footer
    if config.get("include_footer", True):
        for section in doc.sections:
            footer = section.footer
            footer.paragraphs[0].text = f"Entropy.sa — Confidential"

    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _docx_to_pdf(docx_bytes: bytes) -> bytes:
    """Convert DOCX to PDF. Uses LibreOffice headless if available, else returns docx."""
    import subprocess
    import tempfile
    import os

    with tempfile.TemporaryDirectory() as tmpdir:
        docx_path = os.path.join(tmpdir, "proposal.docx")
        with open(docx_path, "wb") as f:
            f.write(docx_bytes)

        try:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, docx_path],
                capture_output=True,
                timeout=60,
            )
            if result.returncode == 0:
                pdf_path = docx_path.replace(".docx", ".pdf")
                if os.path.exists(pdf_path):
                    with open(pdf_path, "rb") as f:
                        return f.read()
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass

    # Fallback: return DOCX bytes (PDF conversion not available)
    return docx_bytes


async def _generate_pptx(proposal, config: dict, language: str) -> bytes:
    """Generate a PPTX presentation from proposal sections using python-pptx."""
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]  # title slide

    PRIMARY = RGBColor(0x48, 0x65, 0x81)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x1E, 0x29, 0x3B)

    def _add_bg(slide, color: RGBColor):
        from pptx.util import Emu
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_text_box(slide, text, left, top, width, height, font_size, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

    # ── Cover slide ─────────────────────────────────────────────────────────
    if config.get("include_cover", True):
        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, PRIMARY)
        _add_text_box(slide, "Entropy.sa", 1, 2.5, 11, 1.2, 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, "Technical Proposal", 1, 3.9, 11, 0.8, 22, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Section slides ───────────────────────────────────────────────────────
    for idx, section in enumerate(proposal.sections, 1):
        title_text = section.title_en if language == "en" else (section.title_ar or section.title_en or f"Section {idx}")
        content_text = section.content_en if language == "en" else (section.content_ar or section.content_en or "")

        if config.get("include_section_numbers", True):
            title_text = f"{idx}. {title_text}"

        slide = prs.slides.add_slide(blank_layout)
        _add_bg(slide, WHITE)

        # Title bar
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(13.33), Inches(1.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()

        _add_text_box(slide, title_text, 0.3, 0.1, 12.5, 0.9, 20, bold=True, color=WHITE)

        # Content (truncate to ~800 chars to fit slide)
        body = (content_text or "").strip()[:800]
        if body:
            _add_text_box(slide, body, 0.5, 1.3, 12.3, 5.8, 14, color=DARK)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


# Public alias used by the BackgroundTasks fallback in proposal.py when Celery is unavailable
